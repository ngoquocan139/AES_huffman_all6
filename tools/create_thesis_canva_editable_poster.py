# -*- coding: utf-8 -*-
"""Create an editable A1 poster PPTX for Canva.

All figures are extracted directly from:
H:\\Academic\\senior_project\\DATN\\docs\\Graduation_Thesis_an_tan.docx

The slide is A1 landscape and uses editable text boxes/shapes. Body text is
kept at 32 pt and section titles at 36 pt as requested.
"""

from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(r"H:\Academic\senior_project\DATN")
DOCS = ROOT / "docs"
DOCX = DOCS / "Graduation_Thesis_an_tan.docx"
OUT_DIR = DOCS / "poster_a1_canva"
MEDIA_DIR = OUT_DIR / "thesis_an_tan_direct_media"
OUT_PPTX = OUT_DIR / "rv32i_huffman_aes_a1_canva_editable_from_thesis.pptx"

BLUE = RGBColor(0, 58, 160)
DARK_BLUE = RGBColor(0, 35, 115)
LIGHT_BLUE = RGBColor(231, 241, 255)
PALE_BLUE = RGBColor(248, 251, 255)
ORANGE = RGBColor(242, 138, 32)
GREEN = RGBColor(34, 145, 92)
RED = RGBColor(210, 72, 72)
GRID = RGBColor(190, 216, 240)
BLACK = RGBColor(25, 31, 42)
WHITE = RGBColor(255, 255, 255)


def extract_docx_media() -> None:
    MEDIA_DIR.mkdir(parents=True, exist_ok=True)
    with ZipFile(DOCX) as zf:
        for name in zf.namelist():
            if name.startswith("word/media/"):
                target = MEDIA_DIR / Path(name).name
                if not target.exists():
                    target.write_bytes(zf.read(name))


def add_textbox(slide, x, y, w, h, text, size=32, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, size=32):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(9)
        p._p.get_or_add_pPr().set("marL", "457200")
        p._p.get_or_add_pPr().set("indent", "-228600")
        p._p.get_or_add_pPr().set("algn", "l")
        # PowerPoint will render the bullet after import; this keeps it editable.
        p._p.get_or_add_pPr().insert(0, p._p._element.makeelement("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar", {"char": "•"}))
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=BLUE, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.8)
    return shp


def add_section(slide, x, y, w, h, label, header_w=None):
    add_rect(slide, x, y, w, h, fill=WHITE, line=RGBColor(40, 145, 225), radius=True)
    hw = header_w or min(w * 0.72, w - 0.2)
    hdr = add_rect(slide, x, y, hw, 0.58, fill=BLUE, line=BLUE, radius=True)
    hdr.adjustments[0] = 0.12
    add_textbox(slide, x + 0.18, y + 0.11, hw - 0.3, 0.42, label, size=36, bold=True, color=WHITE)
    return y + 0.74


def add_picture_fit(slide, path: Path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pic_w = iw * scale
    pic_h = ih * scale
    left = x + (w - pic_w) / 2
    top = y + (h - pic_h) / 2
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(pic_w), height=Inches(pic_h))


def add_stat(slide, x, y, w, h, value, label, color):
    add_rect(slide, x, y, w, h, fill=PALE_BLUE, line=GRID, radius=True)
    add_textbox(slide, x + 0.05, y + 0.12, w - 0.1, 0.38, value, size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.08, y + 0.62, w - 0.16, h - 0.65, label, size=32, bold=True, color=BLACK, align=PP_ALIGN.CENTER)


def add_flow_node(slide, x, y, w, h, text, fill=LIGHT_BLUE, line=BLUE):
    shp = add_rect(slide, x, y, w, h, fill=fill, line=line, radius=True)
    shp.adjustments[0] = 0.12
    add_textbox(slide, x + 0.06, y + 0.13, w - 0.12, h - 0.18, text, size=32, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=BLUE):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(2.2)
    line.line.end_arrowhead = True
    return line


def main() -> None:
    extract_docx_media()

    prs = Presentation()
    prs.slide_width = Inches(33.11)
    prs.slide_height = Inches(23.39)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0.05, 0.05, 33.01, 23.29, fill=WHITE, line=RGBColor(45, 45, 45), radius=False)

    # Header
    add_picture_fit(slide, MEDIA_DIR / "image1.jpeg", 0.35, 0.25, 2.0, 1.9)
    add_textbox(slide, 2.45, 0.35, 5.0, 0.25, "TRƯỜNG ĐẠI HỌC SƯ PHẠM KỸ THUẬT", size=32, bold=True, color=DARK_BLUE)
    add_textbox(slide, 2.45, 0.72, 4.2, 0.25, "THÀNH PHỐ HỒ CHÍ MINH", size=32, bold=True, color=DARK_BLUE)
    add_textbox(slide, 2.45, 1.24, 3.0, 0.25, "KHOA ĐIỆN - ĐIỆN TỬ", size=32, bold=True, color=DARK_BLUE)
    add_textbox(slide, 2.45, 1.60, 5.5, 0.25, "BỘ MÔN KỸ THUẬT MÁY TÍNH - VIỄN THÔNG", size=32, bold=True, color=DARK_BLUE)
    add_rect(slide, 8.60, 0.18, 0.03, 1.95, fill=DARK_BLUE, line=DARK_BLUE)
    add_textbox(slide, 9.0, 0.24, 16.2, 0.42, "DESIGN OF A RISC-V RV32I SYSTEM INTEGRATING", size=36, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 9.0, 0.72, 16.2, 0.42, "HUFFMAN COMPRESSION AND AES-128", size=36, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 9.0, 1.20, 16.2, 0.42, "FOR SECURE DATA STORAGE", size=36, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 27.4, 0.28, 1.1, 0.35, "GVHD:", size=32, bold=True, color=DARK_BLUE)
    add_textbox(slide, 28.4, 0.28, 3.1, 0.35, "PGS.TS. Đỗ Duy Tân", size=32, bold=True, color=BLACK)
    add_textbox(slide, 27.4, 0.95, 1.1, 0.35, "SVTH:", size=32, bold=True, color=DARK_BLUE)
    add_textbox(slide, 28.4, 0.95, 4.1, 0.35, "Ngô Quốc An - 21161226", size=32, bold=True, color=BLACK)
    add_textbox(slide, 28.4, 1.33, 4.1, 0.35, "Bùi Ngọc Duy Tân - 21161266", size=32, bold=True, color=BLACK)

    # Geometry
    x1, w1 = 0.35, 9.35
    x2, w2 = 10.0, 13.5
    x3, w3 = 23.8, 8.95
    y0 = 2.45
    top_h = 6.0
    mid_y = 8.75
    mid_h = 8.6
    bot_y = 17.65
    bot_h = 5.35

    # 1
    cy = add_section(slide, x1, y0, w1, top_h, "1. ABSTRACT / PROBLEM", header_w=5.6)
    add_bullets(slide, x1 + 0.35, cy + 0.25, w1 - 0.65, top_h - 1.0, [
        "Design an RV32I-controlled secure-storage system combining dynamic Huffman compression with AES-128-CBC encryption.",
        "Target flow is compression-to-storage, not data transmission.",
        "TX compresses plaintext, packs 128-bit transport words, encrypts them, and stores ciphertext in DMEM.",
        "RX decrypts and decompresses selected stored records to recover plaintext exactly.",
    ], size=32)

    # 2
    cy = add_section(slide, x1, mid_y, w1, 4.35, "2. OBJECTIVES & SCOPE", header_w=5.2)
    add_bullets(slide, x1 + 0.35, cy + 0.15, w1 - 0.65, 3.35, [
        "Use RV32I as the control and storage-management processor.",
        "Use DMA and RTL accelerators for compression, encryption, decryption, and decompression.",
        "Manage file identifiers, metadata records, plaintext/ciphertext lengths, and IV values in firmware.",
    ], size=32)

    # 7
    cy = add_section(slide, x1, mid_y + 4.65, w1, mid_h - 4.65, "7. FPGA IMPLEMENTATION", header_w=5.0)
    add_stat(slide, x1 + 0.35, cy + 0.10, 1.95, 1.05, "300", "MHz input clock", BLUE)
    add_stat(slide, x1 + 2.55, cy + 0.10, 1.95, 1.05, "50", "MHz SoC clock", BLUE)
    add_stat(slide, x1 + 4.75, cy + 0.10, 2.05, 1.05, "+10.862", "ns WNS", GREEN)
    add_stat(slide, x1 + 7.05, cy + 0.10, 1.95, 1.05, "0.788W", "power", ORANGE)
    add_bullets(slide, x1 + 0.35, cy + 1.40, w1 - 0.65, 1.7, [
        "Full FPGA SoC uses 36,649 LUTs, 20,017 FFs, and 11 BRAMs.",
        "Timing constraints are met on the ZCU102 implementation.",
    ], size=32)

    # 3
    cy = add_section(slide, x2, y0, w2, top_h, "3. SYSTEM ARCHITECTURE", header_w=6.0)
    add_flow_node(slide, x2 + 0.55, cy + 0.20, 2.15, 0.75, "RV32I CPU")
    add_flow_node(slide, x2 + 3.30, cy + 0.20, 2.45, 0.75, "MMIO/APB Bridge")
    add_flow_node(slide, x2 + 6.35, cy + 0.20, 2.70, 0.75, "DMA Register + IV")
    add_arrow(slide, x2 + 2.70, cy + 0.58, x2 + 3.30, cy + 0.58)
    add_arrow(slide, x2 + 5.75, cy + 0.58, x2 + 6.35, cy + 0.58)
    add_flow_node(slide, x2 + 0.75, cy + 2.05, 2.75, 1.0, "DMEM\ninput / metadata / output", fill=RGBColor(255, 248, 230), line=ORANGE)
    add_flow_node(slide, x2 + 4.15, cy + 1.75, 2.10, 0.75, "TX DMA", fill=RGBColor(232, 250, 242), line=GREEN)
    add_flow_node(slide, x2 + 4.15, cy + 3.10, 2.10, 0.75, "RX DMA", fill=RGBColor(232, 250, 242), line=GREEN)
    add_flow_node(slide, x2 + 7.0, cy + 1.72, 2.7, 0.82, "TX Accelerator\nHuffman + AES-CBC")
    add_flow_node(slide, x2 + 7.0, cy + 3.05, 2.7, 0.82, "RX Accelerator\nAES-CBC + Huffman")
    add_flow_node(slide, x2 + 10.45, cy + 2.35, 2.3, 0.95, "Metadata\nfile_id / length / IV", fill=RGBColor(246, 238, 255), line=RGBColor(120, 80, 200))
    add_arrow(slide, x2 + 3.50, cy + 2.55, x2 + 4.15, cy + 2.12, GREEN)
    add_arrow(slide, x2 + 6.25, cy + 2.12, x2 + 7.0, cy + 2.12, GREEN)
    add_arrow(slide, x2 + 3.50, cy + 2.55, x2 + 4.15, cy + 3.48, GREEN)
    add_arrow(slide, x2 + 6.25, cy + 3.48, x2 + 7.0, cy + 3.48, GREEN)
    add_bullets(slide, x2 + 0.45, cy + 4.35, w2 - 0.9, 1.0, [
        "Control plane: RV32I configures DMA registers, writes IV, starts TX/RX, and polls status.",
        "Data plane: DMA and accelerators move data between DMEM, Huffman, and AES blocks.",
    ], size=32)

    # 4
    cy = add_section(slide, x2, mid_y, w2, mid_h, "4. TX ACCELERATOR DESIGN", header_w=6.5)
    add_rect(slide, x2 + 0.35, cy + 0.05, w2 - 0.7, 5.25, fill=PALE_BLUE, line=GRID, radius=True)
    add_picture_fit(slide, MEDIA_DIR / "image12.png", x2 + 0.70, cy + 0.20, w2 - 1.40, 4.95)
    add_bullets(slide, x2 + 0.45, cy + 5.45, w2 - 0.9, 2.05, [
        "The current RTL builds one Huffman table for the whole file.",
        "The first emitted block carries the full table; later blocks reuse the same codebook with compact headers.",
        "Transport words are encrypted by AES-128-CBC or bypassed in COMPRESS_ONLY mode.",
    ], size=32)

    # 5
    cy = add_section(slide, x3, y0, w3, top_h, "5. VERIFICATION RESULTS", header_w=5.4)
    add_flow_node(slide, x3 + 0.35, cy + 0.25, 1.75, 0.65, "Firmware C")
    add_flow_node(slide, x3 + 2.55, cy + 0.25, 1.90, 0.65, "Questa RTL\nSimulation")
    add_flow_node(slide, x3 + 4.95, cy + 0.25, 1.75, 0.65, "DMEM Dump\nCompare")
    add_flow_node(slide, x3 + 7.05, cy + 0.25, 1.55, 0.65, "PASS/FAIL")
    add_arrow(slide, x3 + 2.10, cy + 0.58, x3 + 2.55, cy + 0.58)
    add_arrow(slide, x3 + 4.45, cy + 0.58, x3 + 4.95, cy + 0.58)
    add_arrow(slide, x3 + 6.70, cy + 0.58, x3 + 7.05, cy + 0.58)
    add_stat(slide, x3 + 0.35, cy + 1.35, 3.85, 1.15, "18/0", "TX-RX loopback", GREEN)
    add_stat(slide, x3 + 4.50, cy + 1.35, 3.85, 1.15, "15/0", "Huffman-only input1", GREEN)
    add_stat(slide, x3 + 0.35, cy + 2.75, 3.85, 1.15, "22/0", "Storage table", GREEN)
    add_stat(slide, x3 + 4.50, cy + 2.75, 3.85, 1.15, "34/34", "Regression baseline", GREEN)
    add_bullets(slide, x3 + 0.35, cy + 4.25, w3 - 0.7, 1.15, [
        "RX restored plaintext matches the original source data.",
        "Storage-table test stores multiple records and selects one by file_id.",
    ], size=32)

    # 6
    cy = add_section(slide, x3, mid_y, w3, mid_h, "6. CONCLUSION & FUTURE WORK", header_w=5.9)
    add_bullets(slide, x3 + 0.35, cy + 0.25, w3 - 0.7, 3.2, [
        "The thesis presented an RV32I system integrating dynamic Huffman compression and AES-128-CBC encryption for secure data storage.",
        "Secure-storage firmware manages metadata records, file identifiers, plaintext/ciphertext lengths, and IV values.",
        "The design is functionally verified, FPGA-implementable, and suitable for compact embedded secure data storage.",
    ], size=32)
    add_rect(slide, x3 + 0.35, cy + 5.35, w3 - 0.70, 2.35, fill=PALE_BLUE, line=GRID, radius=True)
    add_textbox(slide, x3 + 0.65, cy + 5.55, 2.6, 0.35, "Future work", size=36, bold=True, color=DARK_BLUE)
    add_bullets(slide, x3 + 0.65, cy + 6.02, w3 - 1.3, 1.45, [
        "Use stronger random or nonce-management mechanism for IV generation.",
        "Add authentication tag to detect malicious ciphertext modification.",
        "Optimize RX decoder resources and extend the number of storage records.",
    ], size=32)

    # 8
    perf_w = 17.45
    cy = add_section(slide, x1, bot_y, perf_w, bot_h, "8. PERFORMANCE COMPARISON", header_w=5.6)
    add_stat(slide, x1 + 0.35, cy + 0.25, 2.35, 1.1, "145.5", "MB/s AES @100MHz", ORANGE)
    add_stat(slide, x1 + 3.0, cy + 0.25, 2.35, 1.1, "0.078", "B/cycle TX path", ORANGE)
    add_stat(slide, x1 + 5.65, cy + 0.25, 2.35, 1.1, "65.50%", "input1 storage saving", GREEN)
    add_stat(slide, x1 + 8.30, cy + 0.25, 2.35, 1.1, "70.13%", "raw ECG storage saving", GREEN)
    add_rect(slide, x1 + 11.0, cy + 0.18, perf_w - 11.35, 3.55, fill=WHITE, line=GRID, radius=True)
    add_picture_fit(slide, MEDIA_DIR / "image59.png", x1 + 11.2, cy + 0.35, perf_w - 11.75, 3.2)
    add_bullets(slide, x1 + 0.35, bot_y + bot_h - 0.75, perf_w - 0.70, 0.55, [
        "The proposed design balances hardware speed, compression efficiency, and secure-storage functionality.",
    ], size=32)

    # 9 thesis sources, no external material added
    cy = add_section(slide, x1 + perf_w + 0.34, bot_y, 33.11 - (x1 + perf_w + 0.34) - 0.35, bot_h, "9. THESIS SOURCES", header_w=4.9)
    add_bullets(slide, x1 + perf_w + 0.70, cy + 0.25, 33.11 - (x1 + perf_w + 0.34) - 1.05, 3.6, [
        "Graduation_Thesis_an_tan.docx: Chapter 1 overview and objectives.",
        "Chapter 3: RV32I secure-storage SoC architecture and TX/RX datapaths.",
        "Chapter 4: verification, performance comparison, and FPGA implementation results.",
        "Chapter 5: conclusion, limitations, and future work.",
    ], size=32)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(OUT_PPTX)


if __name__ == "__main__":
    main()
